#!/usr/bin/env python3
"""Generate an evidence bundle for a PPTX deck.

This script intentionally does NOT generate grounded.md.
It only writes the extraction bundle used later by the agent.
"""

from __future__ import annotations

import argparse
import csv
import json
import math
import re
from dataclasses import dataclass, asdict
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# -----------------------------
# Helpers
# -----------------------------


CONTROL_CHAR_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]")
WHITESPACE_RE = re.compile(r"[ \t]+")


def slugify(text: str) -> str:
    text = text.strip().lower()
    text = re.sub(r"[^a-z0-9]+", "-", text)
    text = re.sub(r"-+", "-", text).strip("-")
    return text or "deck"


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def safe_text(text: Optional[str]) -> str:
    if not text:
        return ""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = CONTROL_CHAR_RE.sub("", text)
    lines = [WHITESPACE_RE.sub(" ", line).strip() for line in text.split("\n")]
    lines = [line for line in lines if line]
    return "\n".join(lines).strip()


def normalize_text_for_compare(text: str) -> str:
    text = safe_text(text).lower()
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def count_words(text: str) -> int:
    text = safe_text(text)
    if not text:
        return 0
    return len(re.findall(r"\S+", text))


def count_lines(text: str) -> int:
    text = safe_text(text)
    if not text:
        return 0
    return len(text.splitlines())


def is_mostly_symbolic(text: str) -> bool:
    compact = re.sub(r"\s+", "", safe_text(text))
    if not compact:
        return True
    symbolic_chars = sum(1 for ch in compact if not ch.isalnum())
    return symbolic_chars / max(len(compact), 1) >= 0.6


def shape_position(shape: Any) -> Dict[str, Optional[int]]:
    return {
        "left": getattr(shape, "left", None),
        "top": getattr(shape, "top", None),
        "width": getattr(shape, "width", None),
        "height": getattr(shape, "height", None),
    }


def get_placeholder_kind(shape: Any) -> Optional[str]:
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        return str(shape.placeholder_format.type)
    except Exception:
        return "placeholder"


def paragraph_to_item(paragraph: Any) -> Optional[Dict[str, Any]]:
    text = safe_text(getattr(paragraph, "text", ""))
    if not text:
        return None
    return {
        "level": int(getattr(paragraph, "level", 0) or 0),
        "text": text,
    }


def get_max_font_pt(shape: Any) -> Optional[float]:
    if not getattr(shape, "has_text_frame", False):
        return None
    best: Optional[float] = None
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                size = getattr(getattr(run, "font", None), "size", None)
                if size is None:
                    continue
                value = float(size.pt)
                if best is None or value > best:
                    best = value
    except Exception:
        return best
    return best


def get_shape_metrics(shape: Any, slide_width: int, slide_height: int) -> Dict[str, Optional[float]]:
    left = getattr(shape, "left", None)
    top = getattr(shape, "top", None)
    width = getattr(shape, "width", None)
    height = getattr(shape, "height", None)

    return {
        "left_ratio": (float(left) / slide_width) if left is not None and slide_width else None,
        "top_ratio": (float(top) / slide_height) if top is not None and slide_height else None,
        "width_ratio": (float(width) / slide_width) if width is not None and slide_width else None,
        "height_ratio": (float(height) / slide_height) if height is not None and slide_height else None,
        "area_ratio": (
            float(width * height) / float(slide_width * slide_height)
            if width is not None and height is not None and slide_width and slide_height
            else None
        ),
    }


@dataclass
class AssetRecord:
    asset_id: str
    asset_type: str
    slide_no: int
    source_shape_id: Optional[int]
    source_shape_name: Optional[str]
    rel_path: Optional[str]
    extra: Dict[str, Any]
    extraction_status: str


# -----------------------------
# Extraction primitives
# -----------------------------


def iter_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        yield shape
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            for sub in iter_shapes(shape.shapes):
                yield sub


def detect_title(slide: Any, slide_width: int, slide_height: int) -> Tuple[str, Optional[int], str]:
    try:
        title_shape = slide.shapes.title
    except Exception:
        title_shape = None

    if title_shape is not None:
        text = safe_text(getattr(title_shape, "text", ""))
        if text:
            return text, getattr(title_shape, "shape_id", None), "slide.shapes.title"

    best_text = ""
    best_shape_id: Optional[int] = None
    best_score = -10**9

    for shape in iter_shapes(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue

        text = safe_text(getattr(shape, "text", ""))
        if not text:
            continue

        placeholder_kind = (get_placeholder_kind(shape) or "").upper()
        word_count = count_words(text)
        line_count = count_lines(text)
        metrics = get_shape_metrics(shape, slide_width, slide_height)
        top_ratio = metrics.get("top_ratio")
        width_ratio = metrics.get("width_ratio")
        max_font_pt = get_max_font_pt(shape) or 0.0

        score = 0.0

        if "TITLE" in placeholder_kind:
            score += 1000
        elif "CENTER_TITLE" in placeholder_kind:
            score += 950
        elif "SUBTITLE" in placeholder_kind:
            score += 200

        if top_ratio is not None:
            if top_ratio <= 0.18:
                score += 220
            elif top_ratio <= 0.28:
                score += 120
            elif top_ratio >= 0.45:
                score -= 120

        if max_font_pt:
            score += min(max_font_pt, 40.0) * 8.0

        if width_ratio is not None and 0.18 <= width_ratio <= 0.95:
            score += 40

        # Prefer concise titles; strongly penalize long正文.
        if word_count <= 8:
            score += 140
        elif word_count <= 16:
            score += 60
        elif word_count <= 28:
            score += 10
        else:
            score -= min((word_count - 28) * 10, 300)

        if line_count == 1:
            score += 80
        elif line_count == 2:
            score += 20
        else:
            score -= min((line_count - 2) * 25, 150)

        if "\n" in text and word_count > 18:
            score -= 120

        if is_mostly_symbolic(text):
            score -= 180

        if score > best_score:
            best_score = score
            best_text = text.split("\n", 1)[0].strip()
            best_shape_id = getattr(shape, "shape_id", None)

    return best_text, best_shape_id, "heuristic"


def classify_text_shape(
    shape: Any,
    slide_title: str,
    title_shape_id: Optional[int],
    slide_width: int,
    slide_height: int,
) -> Tuple[Optional[Dict[str, Any]], Optional[Dict[str, Any]], Optional[Dict[str, Any]]]:
    if not getattr(shape, "has_text_frame", False):
        return None, None, None

    full_text = safe_text(getattr(shape, "text", ""))
    if not full_text:
        return None, None, None

    shape_id = getattr(shape, "shape_id", None)
    if title_shape_id is not None and shape_id == title_shape_id:
        return None, None, None

    if slide_title and normalize_text_for_compare(full_text) == normalize_text_for_compare(slide_title):
        return None, None, None

    paragraphs: List[Dict[str, Any]] = []
    for p in getattr(shape.text_frame, "paragraphs", []):
        item = paragraph_to_item(p)
        if item:
            paragraphs.append(item)

    if not paragraphs:
        return None, None, None

    placeholder_kind = get_placeholder_kind(shape)
    metrics = get_shape_metrics(shape, slide_width, slide_height)
    word_count = count_words(full_text)
    line_count = count_lines(full_text)
    max_font_pt = get_max_font_pt(shape)

    meta = {
        "shape_id": shape_id,
        "shape_name": getattr(shape, "name", None),
        "placeholder_kind": placeholder_kind,
        "position": shape_position(shape),
        "text_metrics": {
            "word_count": word_count,
            "line_count": line_count,
            "max_font_pt": max_font_pt,
            **metrics,
        },
    }

    any_nested = any((item.get("level", 0) or 0) > 0 for item in paragraphs)
    multi_paragraph = len(paragraphs) >= 2
    looks_like_body_placeholder = bool(placeholder_kind and "BODY" in str(placeholder_kind).upper())
    is_bullet_like = multi_paragraph or any_nested or looks_like_body_placeholder

    if is_bullet_like:
        return {
            **meta,
            "role": "bullet_block",
            "items": paragraphs,
            "raw_text": full_text,
        }, None, None

    compact = re.sub(r"\s+", "", full_text)
    top_ratio = metrics.get("top_ratio")
    area_ratio = metrics.get("area_ratio")
    max_font_pt_val = max_font_pt or 0.0
    placeholder_upper = (placeholder_kind or "").upper()

    is_short = word_count <= 8 and len(compact) <= 36
    is_topish = top_ratio is not None and top_ratio <= 0.26
    is_really_small = area_ratio is not None and area_ratio <= 0.015
    is_label_like = is_short and (is_really_small or len(compact) <= 16)

    if is_short and (
        "SUBTITLE" in placeholder_upper
        or "HEADER" in placeholder_upper
        or (is_topish and max_font_pt_val >= 18)
        or (is_topish and line_count <= 2 and len(compact) <= 24)
    ):
        return None, {
            **meta,
            "role": "section_marker",
            "text": full_text,
        }, None

    if is_mostly_symbolic(full_text) or (len(compact) <= 4 and not re.search(r"[A-Za-z\u4e00-\u9fff]", compact)):
        return None, None, {
            **meta,
            "role": "decorative_text",
            "text": full_text,
        }

    if is_label_like:
        return None, {
            **meta,
            "role": "caption_or_label",
            "text": full_text,
        }, None

    return None, {
        **meta,
        "role": "body_text",
        "text": full_text,
    }, None


def extract_notes(slide: Any) -> str:
    try:
        if not slide.has_notes_slide:
            return ""
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        return safe_text(text_frame.text)
    except Exception:
        return ""


def export_picture(shape: Any, slide_no: int, images_dir: Path) -> Tuple[AssetRecord, Optional[str]]:
    try:
        image = shape.image
        ext = image.ext or "bin"
        asset_id = f"slide_{slide_no:03d}_image_{getattr(shape, 'shape_id', 0):03d}"
        filename = f"{asset_id}.{ext}"
        out_path = images_dir / filename
        out_path.write_bytes(image.blob)

        size = getattr(image, "size", None)
        record = AssetRecord(
            asset_id=asset_id,
            asset_type="image",
            slide_no=slide_no,
            source_shape_id=getattr(shape, "shape_id", None),
            source_shape_name=getattr(shape, "name", None),
            rel_path=str(Path("assets/images") / filename),
            extra={
                "content_type": getattr(image, "content_type", None),
                "ext": ext,
                "size_px": {
                    "width": size[0] if size else None,
                    "height": size[1] if size else None,
                },
                "position": shape_position(shape),
            },
            extraction_status="ok",
        )
        return record, filename
    except Exception as exc:
        record = AssetRecord(
            asset_id=f"slide_{slide_no:03d}_image_error_{getattr(shape, 'shape_id', 0):03d}",
            asset_type="image",
            slide_no=slide_no,
            source_shape_id=getattr(shape, "shape_id", None),
            source_shape_name=getattr(shape, "name", None),
            rel_path=None,
            extra={"error": str(exc), "position": shape_position(shape)},
            extraction_status="failed",
        )
        return record, None


def write_table_markdown(rows: List[List[str]]) -> str:
    if not rows:
        return ""
    header = rows[0]
    sep = ["---"] * len(header)
    md_rows = ["| " + " | ".join(header) + " |", "| " + " | ".join(sep) + " |"]
    for row in rows[1:]:
        padded = row + [""] * max(0, len(header) - len(row))
        md_rows.append("| " + " | ".join(padded[: len(header)]) + " |")
    return "\n".join(md_rows)


def export_table(shape: Any, slide_no: int, tables_dir: Path) -> AssetRecord:
    asset_id = f"slide_{slide_no:03d}_table_{getattr(shape, 'shape_id', 0):03d}"
    md_name = f"{asset_id}.md"
    csv_name = f"{asset_id}.csv"
    md_path = tables_dir / md_name
    csv_path = tables_dir / csv_name

    try:
        table = shape.table
        rows: List[List[str]] = []
        for row in table.rows:
            row_values = []
            for cell in row.cells:
                row_values.append(safe_text(cell.text))
            rows.append(row_values)

        md_path.write_text(write_table_markdown(rows), encoding="utf-8")
        with csv_path.open("w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerows(rows)

        return AssetRecord(
            asset_id=asset_id,
            asset_type="table",
            slide_no=slide_no,
            source_shape_id=getattr(shape, "shape_id", None),
            source_shape_name=getattr(shape, "name", None),
            rel_path=str(Path("assets/tables") / md_name),
            extra={
                "csv_path": str(Path("assets/tables") / csv_name),
                "rows": len(rows),
                "cols": max((len(r) for r in rows), default=0),
                "position": shape_position(shape),
            },
            extraction_status="ok",
        )
    except Exception as exc:
        return AssetRecord(
            asset_id=asset_id,
            asset_type="table",
            slide_no=slide_no,
            source_shape_id=getattr(shape, "shape_id", None),
            source_shape_name=getattr(shape, "name", None),
            rel_path=None,
            extra={"error": str(exc), "position": shape_position(shape)},
            extraction_status="failed",
        )


def try_get_chart_title(chart: Any) -> str:
    try:
        if chart.has_title:
            return safe_text(chart.chart_title.text_frame.text)
    except Exception:
        pass
    return ""


def try_get_chart_categories(chart: Any) -> List[str]:
    candidates: List[str] = []
    try:
        plot = chart.plots[0]
        categories = plot.categories
        for c in categories:
            candidates.append(safe_text(str(c.label)))
        if candidates:
            return candidates
    except Exception:
        pass
    return []


def export_chart(shape: Any, slide_no: int, charts_dir: Path) -> AssetRecord:
    asset_id = f"slide_{slide_no:03d}_chart_{getattr(shape, 'shape_id', 0):03d}"
    json_name = f"{asset_id}.json"
    csv_name = f"{asset_id}.csv"
    json_path = charts_dir / json_name
    csv_path = charts_dir / csv_name

    try:
        chart = shape.chart
        categories = try_get_chart_categories(chart)
        series_payload = []
        flat_rows: List[List[Any]] = [["series", "category", "value"]]

        for series in chart.series:
            try:
                values = list(series.values)
            except Exception:
                values = []
            try:
                series_name = safe_text(str(series.name))
            except Exception:
                series_name = ""

            series_payload.append({
                "name": series_name,
                "values": values,
            })
            for idx, value in enumerate(values):
                category = categories[idx] if idx < len(categories) else str(idx)
                flat_rows.append([series_name, category, value])

        payload = {
            "chart_type": safe_text(str(getattr(chart, "chart_type", ""))),
            "title": try_get_chart_title(chart),
            "categories": categories,
            "series": series_payload,
            "position": shape_position(shape),
        }

        json_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
        with csv_path.open("w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerows(flat_rows)

        return AssetRecord(
            asset_id=asset_id,
            asset_type="chart",
            slide_no=slide_no,
            source_shape_id=getattr(shape, "shape_id", None),
            source_shape_name=getattr(shape, "name", None),
            rel_path=str(Path("assets/charts") / json_name),
            extra={
                "csv_path": str(Path("assets/charts") / csv_name),
                "series_count": len(series_payload),
                "category_count": len(categories),
                "position": shape_position(shape),
            },
            extraction_status="ok",
        )
    except Exception as exc:
        return AssetRecord(
            asset_id=asset_id,
            asset_type="chart",
            slide_no=slide_no,
            source_shape_id=getattr(shape, "shape_id", None),
            source_shape_name=getattr(shape, "name", None),
            rel_path=None,
            extra={"error": str(exc), "position": shape_position(shape)},
            extraction_status="failed",
        )


# -----------------------------
# Main extraction
# -----------------------------


def build_bundle(input_pptx: Path, output_root: Path) -> Path:
    prs = Presentation(str(input_pptx))

    deck_id = slugify(input_pptx.stem)

    # Generate GROUND_ID: pptx-deck_id_timestamp (Beijing time)
    tz_beijing = timezone(timedelta(hours=8))
    timestamp = datetime.now(tz_beijing).strftime("%Y%m%d%H%M%S")
    ground_id = f"pptx-{deck_id}_{timestamp}"
    bundle_dir = output_root / ground_id
    ensure_dir(bundle_dir)

    # Write ground_id.txt for downstream stages to reuse
    (bundle_dir / "ground_id.txt").write_text(ground_id + "\n", encoding="utf-8")

    assets_dir = bundle_dir / "assets"
    images_dir = assets_dir / "images"
    tables_dir = assets_dir / "tables"
    charts_dir = assets_dir / "charts"
    for d in (assets_dir, images_dir, tables_dir, charts_dir):
        ensure_dir(d)

    slide_width = int(getattr(prs, "slide_width", 0) or 0)
    slide_height = int(getattr(prs, "slide_height", 0) or 0)

    title_from_props = safe_text(getattr(getattr(prs, "core_properties", None), "title", ""))
    parser_name = "python-pptx"
    extracted_at = datetime.now(timezone.utc).isoformat()

    slide_records: List[Dict[str, Any]] = []
    asset_records: List[Dict[str, Any]] = []

    note_count = 0
    image_count = 0
    table_count = 0
    chart_count = 0

    for idx, slide in enumerate(prs.slides, start=1):
        slide_title, title_shape_id, title_source = detect_title(slide, slide_width, slide_height)
        notes_text = extract_notes(slide)
        if notes_text:
            note_count += 1

        bullet_blocks: List[Dict[str, Any]] = []
        text_blocks: List[Dict[str, Any]] = []
        section_markers: List[Dict[str, Any]] = []
        decorative_texts: List[Dict[str, Any]] = []
        slide_assets: List[Dict[str, Any]] = []

        shape_stats = {
            "total_shapes": 0,
            "text_shapes": 0,
            "image_shapes": 0,
            "table_shapes": 0,
            "chart_shapes": 0,
        }

        for shape in iter_shapes(slide.shapes):
            shape_stats["total_shapes"] += 1

            if getattr(shape, "has_text_frame", False):
                shape_stats["text_shapes"] += 1
                bullet_block, text_block, decorative_text = classify_text_shape(
                    shape,
                    slide_title,
                    title_shape_id,
                    slide_width,
                    slide_height,
                )
                if bullet_block:
                    bullet_blocks.append(bullet_block)
                if text_block:
                    if text_block.get("role") == "section_marker":
                        section_markers.append(text_block)
                    else:
                        text_blocks.append(text_block)
                if decorative_text:
                    decorative_texts.append(decorative_text)

            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                shape_stats["image_shapes"] += 1
                image_count += 1
                asset, _ = export_picture(shape, idx, images_dir)
                asset_records.append(asdict(asset))
                slide_assets.append({
                    "asset_id": asset.asset_id,
                    "asset_type": asset.asset_type,
                    "rel_path": asset.rel_path,
                    "status": asset.extraction_status,
                })
                continue

            has_table = bool(getattr(shape, "has_table", False))
            has_chart = bool(getattr(shape, "has_chart", False))

            if has_table:
                shape_stats["table_shapes"] += 1
                table_count += 1
                asset = export_table(shape, idx, tables_dir)
                asset_records.append(asdict(asset))
                slide_assets.append({
                    "asset_id": asset.asset_id,
                    "asset_type": asset.asset_type,
                    "rel_path": asset.rel_path,
                    "status": asset.extraction_status,
                })

            if has_chart:
                shape_stats["chart_shapes"] += 1
                chart_count += 1
                asset = export_chart(shape, idx, charts_dir)
                asset_records.append(asdict(asset))
                slide_assets.append({
                    "asset_id": asset.asset_id,
                    "asset_type": asset.asset_type,
                    "rel_path": asset.rel_path,
                    "status": asset.extraction_status,
                })

        slide_records.append({
            "slide_no": idx,
            "slide_id": f"slide_{idx:03d}",
            "layout_name": safe_text(getattr(getattr(slide, "slide_layout", None), "name", "")),
            "title": slide_title,
            "title_meta": {
                "shape_id": title_shape_id,
                "source": title_source,
            },
            "section_markers": section_markers,
            "bullet_blocks": bullet_blocks,
            "text_blocks": text_blocks,
            "decorative_texts": decorative_texts,
            "notes_text": notes_text,
            "assets": slide_assets,
            "shape_stats": shape_stats,
        })

    extracted_meta = {
        "source_file": str(input_pptx),
        "source_type": "pptx",
        "deck_id": deck_id,
        "bundle_dir": str(bundle_dir),
        "parser": parser_name,
        "extracted_at": extracted_at,
        "slide_count": len(slide_records),
        "slide_size": {
            "width": slide_width,
            "height": slide_height,
        },
        "deck_title": title_from_props or (slide_records[0]["title"] if slide_records else ""),
        "note_count": note_count,
        "image_count": image_count,
        "table_count": table_count,
        "chart_count": chart_count,
        "scope": {
            "video_audio": "out_of_scope",
            "ocr": "out_of_scope",
            "slide_preview": "not_generated",
        },
    }

    slide_index = {
        "source_file": str(input_pptx),
        "deck_id": deck_id,
        "slide_count": len(slide_records),
        "slides": slide_records,
    }

    asset_index = {
        "source_file": str(input_pptx),
        "deck_id": deck_id,
        "asset_count": len(asset_records),
        "assets": asset_records,
    }

    (bundle_dir / "extracted_meta.json").write_text(
        json.dumps(extracted_meta, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    (bundle_dir / "slide_index.json").write_text(
        json.dumps(slide_index, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    (bundle_dir / "asset_index.json").write_text(
        json.dumps(asset_index, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    (bundle_dir / "extracted.md").write_text(
        build_extracted_md(extracted_meta, slide_records), encoding="utf-8"
    )

    return bundle_dir


# -----------------------------
# Markdown synthesis
# -----------------------------


def build_asset_ref(asset: Dict[str, Any]) -> str:
    lines = [
        "[AssetRef]",
        f"id: {asset.get('asset_id', '')}",
        f"type: {asset.get('asset_type', '')}",
        f"status: {asset.get('status', '')}",
    ]
    if asset.get("rel_path"):
        lines.append(f"path: {asset['rel_path']}")
    lines.append("[/AssetRef]")
    return "\n".join(lines)


def emit_role_group(out: List[str], heading: str, blocks: List[Dict[str, Any]]) -> None:
    out.append(f"### {heading}")
    if blocks:
        for block in blocks:
            role = block.get("role")
            role_suffix = f" role={role}" if role else ""
            out.append(f"- shape_id={block.get('shape_id')} shape_name={block.get('shape_name')}{role_suffix}")
            text = safe_text(block.get("text", ""))
            if text:
                for line in text.splitlines():
                    out.append(f"  - {line}")
    else:
        out.append("- none")
    out.append("")


def build_extracted_md(meta: Dict[str, Any], slides: List[Dict[str, Any]]) -> str:
    out: List[str] = []
    out.append("# PPTX Extraction Bundle")
    out.append("")
    out.append("## Deck Overview")
    out.append(f"- Source file: {meta.get('source_file', '')}")
    out.append(f"- Deck ID: {meta.get('deck_id', '')}")
    out.append(f"- Slide count: {meta.get('slide_count', 0)}")
    out.append(f"- Notes count: {meta.get('note_count', 0)}")
    out.append(f"- Image count: {meta.get('image_count', 0)}")
    out.append(f"- Table count: {meta.get('table_count', 0)}")
    out.append(f"- Chart count: {meta.get('chart_count', 0)}")
    out.append(f"- Parser: {meta.get('parser', '')}")
    out.append("")

    for slide in slides:
        out.append(f"## Slide {slide['slide_no']}")
        out.append("")
        out.append(f"- Title: {slide.get('title', '')}")
        title_meta = slide.get("title_meta", {})
        out.append(f"- Title source: {title_meta.get('source', '')}")
        out.append(f"- Layout: {slide.get('layout_name', '')}")
        out.append("")

        emit_role_group(out, "Section Markers", slide.get("section_markers", []))

        out.append("### Bullet Blocks")
        if slide.get("bullet_blocks"):
            for block in slide["bullet_blocks"]:
                out.append(
                    f"- shape_id={block.get('shape_id')} shape_name={block.get('shape_name')} role={block.get('role', '')}"
                )
                for item in block.get("items", []):
                    indent = "  " * int(item.get("level", 0))
                    out.append(f"  {indent}- {item.get('text', '')}")
        else:
            out.append("- none")
        out.append("")

        emit_role_group(out, "Other Text Blocks", slide.get("text_blocks", []))
        emit_role_group(out, "Decorative Text", slide.get("decorative_texts", []))

        out.append("### Speaker Notes")
        notes_text = safe_text(slide.get("notes_text", ""))
        if notes_text:
            for line in notes_text.splitlines():
                out.append(f"- {line}")
        else:
            out.append("- none")
        out.append("")

        out.append("### Assets")
        if slide.get("assets"):
            for asset in slide["assets"]:
                out.append(build_asset_ref(asset))
                out.append("")
        else:
            out.append("- none")
            out.append("")

    return "\n".join(out).strip() + "\n"


# -----------------------------
# CLI
# -----------------------------


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Extract a PPTX evidence bundle.")
    parser.add_argument("--input", required=True, help="Path to .pptx input")
    parser.add_argument("--output-root", required=True, help="Parent directory for bundle output")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    input_pptx = Path(args.input).expanduser().resolve()
    output_root = Path(args.output_root).expanduser().resolve()

    if not input_pptx.exists():
        raise FileNotFoundError(f"Input PPTX not found: {input_pptx}")
    if input_pptx.suffix.lower() != ".pptx":
        raise ValueError(f"Expected a .pptx file, got: {input_pptx}")

    ensure_dir(output_root)
    bundle_dir = build_bundle(input_pptx, output_root)
    print(f"[pptx-grounding] bundle written to: {bundle_dir}")
    print("[pptx-grounding] next step: read the bundle and write a real grounded.md")


if __name__ == "__main__":
    main()
